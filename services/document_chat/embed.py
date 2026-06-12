import groq
import base64
import gc
import os
import fitz
import io
import subprocess
import openpyxl

from PIL import Image
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv

from .processing import make_chunks, embed_texts_iter
from .database import create_collection, insert_stream_to_qdrant
from .utilities import clean_text, is_garbage_text, verify_insert
from .processor import setup_qdrant, load_model  # use cloud Qdrant + shared model
load_dotenv()

# Cap the combined OCR image size aggressively. A 30 MP RGB image is ~90 MB
# raw and the LANCZOS resize allocates a second copy — that alone can blow
# Render's 512 MB budget. 4 MP is plenty for OCR/vision text extraction.
MAX_OCR_PIXELS = int(os.getenv("MAX_OCR_PIXELS", str(4_000_000)))

groq_client = groq.Groq(api_key=os.getenv("groq_api_key"))


def analyze_image_with_groq(image: Image.Image) -> str:
    buffer = io.BytesIO()
    # JPEG is far smaller than PNG in memory/payload and fine for OCR.
    rgb = image.convert("RGB") if image.mode != "RGB" else image
    rgb.save(buffer, format="JPEG", quality=80)
    if rgb is not image:
        rgb.close()
    img_base64 = base64.b64encode(buffer.getvalue()).decode("utf-8")
    buffer.close()
    try:
        response = groq_client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_base64}"}},
                    {"type": "text", "text": "Extract all text from this image. If no text, describe what you see in two or more sentences"}
                ]
            }]
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"groq error: {e}")
        return ""
    finally:
        del img_base64
        gc.collect()


def concatenate_images_vertically(images: list[Image.Image], max_pixels=MAX_OCR_PIXELS) -> Image.Image:
    if not images:
        return None

    # Cap the working width so individual huge images don't explode memory.
    MAX_WIDTH = 1600
    max_width = min(max(img.width for img in images), MAX_WIDTH)

    resized = []
    for img in images:
        if img.width != max_width:
            ratio = max_width / img.width
            new = img.resize((max_width, max(1, int(img.height * ratio))), Image.LANCZOS)
            img.close()  # free the original copy ASAP
            img = new
        resized.append(img)

    total_height = sum(img.height for img in resized)
    combined = Image.new('RGB', (max_width, total_height), 'white')
    y = 0
    for img in resized:
        combined.paste(img, (0, y))
        y += img.height
        img.close()  # free each source after pasting
    resized.clear()

    if combined.width * combined.height > max_pixels:
        scale = (max_pixels / (combined.width * combined.height)) ** 0.5
        shrunk = combined.resize(
            (max(1, int(combined.width * scale)), max(1, int(combined.height * scale))),
            Image.LANCZOS,
        )
        combined.close()
        combined = shrunk

    gc.collect()
    return combined


def load_ppt(file_path):
    original_path = file_path
    temp_converted = False
    if file_path.lower().endswith('.ppt'):
        file_dir = os.path.dirname(file_path) or "."
        cmd = ["soffice", "--headless", "--convert-to", "pptx", "--outdir", file_dir, file_path]
        subprocess.run(cmd, check=True, capture_output=True, timeout=60)
        file_path = file_path.replace('.ppt', '.pptx')
        temp_converted = True

    filename = os.path.basename(original_path)
    prs = Presentation(file_path)

    all_slides = []
    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        slide_images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        slide_text.append(line)
            elif shape.shape_type == 13:
                try:
                    slide_images.append(Image.open(io.BytesIO(shape.image.blob)))
                except Exception:
                    pass

        all_slides.append({'num': slide_num, 'text': slide_text, 'images': slide_images})

    pages = []
    for i in range(0, len(all_slides), 3):
        batch = all_slides[i:i + 3]
        batch_images = [img for s in batch for img in s['images']]

        image_text = ""
        if batch_images:
            combined = concatenate_images_vertically(batch_images)
            if combined is not None:
                image_text = analyze_image_with_groq(combined)
                combined.close()

        for s in batch:
            full_text = " ".join(s['text']) + " " + image_text
            cleaned = clean_text(" ".join(full_text.split()))

            if len(cleaned.strip()) >= 3:
                pages.append({
                    "text":     cleaned,
                    "page_num": s['num'] + 1,
                    "source":   filename
                })

        # Free this batch's images before moving on.
        for s in batch:
            for img in s['images']:
                try:
                    img.close()
                except Exception:
                    pass
            s['images'] = []
        gc.collect()

    if temp_converted and os.path.exists(file_path):
        os.remove(file_path)
    return pages


def extract_images_text_from_pdf_page(doc, page):
    image_texts = []
    for img in page.get_images(full=True):
        xref = img[0]
        try:
            base_image  = doc.extract_image(xref)
            image       = Image.open(io.BytesIO(base_image["image"]))
            ocr_text    = analyze_image_with_groq(image)
            if ocr_text.strip():
                image_texts.append(ocr_text.strip())
        except Exception as e:
            print(f"Skipping pdf image error: {e}")
    return " ".join(image_texts)


def load_pdf(file_path):
    filename = os.path.basename(file_path)
    doc      = fitz.open(file_path)
    pages    = []
    skipped  = 0

    for page_num in range(len(doc)):
        page     = doc[page_num]
        raw_text = page.get_text()
        full_text = raw_text + " "
        if is_garbage_text(full_text):
            skipped += 1
            print(f"[PDF] Skipping page {page_num+1} — garbage text detected")
            continue
        cleaned = clean_text(full_text)
        if len(cleaned) < 50:
            skipped += 1
            continue
        pages.append({
            "text":     cleaned,
            "page_num": page_num + 1,
            "source":   filename
        })
    print(f"[PDF] {filename}: {len(pages)} pages loaded, {skipped} skipped")
    doc.close()
    
    return pages


def load_image(file_path):
    filename = os.path.basename(file_path)
    with Image.open(file_path) as image:
        # Downscale very large images before OCR to cap memory.
        if image.width * image.height > MAX_OCR_PIXELS:
            scale = (MAX_OCR_PIXELS / (image.width * image.height)) ** 0.5
            image = image.resize(
                (max(1, int(image.width * scale)), max(1, int(image.height * scale))),
                Image.LANCZOS,
            )
        text = analyze_image_with_groq(image)
    gc.collect()
    cleaned = clean_text(text)
    if len(cleaned) < 5:
        return []
    return [{"text": cleaned, "page_num": 1, "source": filename}]


def load_docx(file_path):
    filename   = os.path.basename(file_path)
    doc        = Document(file_path)
    pages      = []
    para_batch = []
    page_num   = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            para_batch.append(text)

        if len(para_batch) >= 20:
            full_text = " ".join(para_batch)
            cleaned   = clean_text(full_text)
            if not is_garbage_text(full_text) and len(cleaned) >= 50:
                pages.append({"text": cleaned, "page_num": page_num, "source": filename})
                page_num += 1
            para_batch = []

    if para_batch:
        full_text = " ".join(para_batch)
        cleaned   = clean_text(full_text)
        if not is_garbage_text(full_text) and len(cleaned) >= 50:
            pages.append({"text": cleaned, "page_num": page_num, "source": filename})

    return pages


def load_xlsx(file_path):
    filename = os.path.basename(file_path)
    wb       = openpyxl.load_workbook(file_path, data_only=True)
    pages    = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames):
        ws         = wb[sheet_name]
        sheet_rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(
                str(cell) for cell in row
                if cell is not None and str(cell).strip()
            )
            if row_text.strip():
                sheet_rows.append(row_text)

        full_text = " ".join(sheet_rows)
        if len(full_text.strip()) < 10:
            continue
        pages.append({
            "text":     clean_text(full_text),
            "page_num": sheet_num + 1,
            "source":   filename
        })
    return pages


def load_file(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:
        return load_image(file_path)
    elif ext in [".pptx", ".ppt"]:
        return load_ppt(file_path)
    elif ext in [".xlsx", ".xls"]:
        return load_xlsx(file_path)
    elif ext in [".docx", ".doc"]:
        return load_docx(file_path)
    else:
        print(f"Unsupported file type: {ext}")
        return []


def main_pipeline(file_path: str, user_id: str, qdrant_client=None, embedding_model=None,
                  original_filename: str = None):
    """
    user_id is stored in each Qdrant point's payload so queries can be
    filtered by owner (preventing cross-user data leakage).

    original_filename: the real uploaded name. We store THIS as `source`
    so chat (which filters by the original filename) actually matches.
    Without it, the temp file name (e.g. "temp_foo.pdf") leaks into `source`
    and every filtered search returns nothing.
    """
    from .processor import setup_qdrant as _setup, load_model as _load

    client = qdrant_client if qdrant_client is not None else _setup()
    model  = embedding_model if embedding_model is not None else _load()

    # NOTE: the collection is created lazily inside insert_stream_to_qdrant()
    # using the embedder's actual vector dimension, so we don't create it here.

    # file_path is always a single string — wrap it for uniform processing
    file_paths = [file_path] if isinstance(file_path, str) else file_path

    all_pages: list[dict] = []
    for fp in file_paths:
        print(f"processing: {fp}")
        pages = load_file(fp)
        for page in pages:
            # Force the real upload name as the source (fixes chat filter mismatch)
            if original_filename:
                page["source"] = original_filename
            if user_id:
                page["user_id"] = user_id
        all_pages.extend(pages)
        gc.collect()

    chunks = make_chunks(all_pages)
    if user_id:
        for chunk in chunks:
            chunk["user_id"] = user_id

    # Free the page text now that chunks are built.
    all_pages.clear()
    gc.collect()

    # Stream embeddings -> Qdrant in small batches so we never hold all
    # vectors in RAM at once (keeps peak memory under Render's 512 MB cap).
    texts = [c["chunk_text"] for c in chunks]
    vector_iter = embed_texts_iter(texts, model)
    inserted = insert_stream_to_qdrant(chunks, vector_iter, client)

    del texts, vector_iter, chunks
    try:
        from services.mem import release_memory
        release_memory()
    except Exception:
        gc.collect()

    print(f"[Pipeline] Done. Inserted {inserted} chunks for user {user_id}.")
    return client, model


if __name__ == "__main__":
    files = []
    client, model = main_pipeline(files)